import streamlit as st
import pandas as pd
import plotly.express as px
import os
import re
import requests
from pathlib import Path
from datetime import datetime
import fitz  # PyMuPDF
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize

# --------------------------------------------------
# NLTK SETUP
# --------------------------------------------------
nltk.download("punkt")

# --------------------------------------------------
# PAGE CONFIG
# --------------------------------------------------
st.set_page_config(
    page_title="Faber Nexus | Consulting Knowledge OS",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --------------------------------------------------
# CSS
# --------------------------------------------------
st.markdown("""
<style>
.stApp { background-color: #f8fafc; }
h1, h2, h3 { color: #1e3a5f; }
.stButton>button {
    background-color: #208C8D;
    color: white;
    border-radius: 8px;
    font-weight: 600;
}
.badge {
    padding: 4px 10px;
    border-radius: 12px;
    font-size: 12px;
    font-weight: 600;
}
.verified { background-color: #e6f9ee; color: #067647; }
.indicative { background-color: #fff4e5; color: #92400e; }
</style>
""", unsafe_allow_html=True)

# --------------------------------------------------
# MASTER LISTS
# --------------------------------------------------
INDUSTRIES = [
    "Automotive", "Pharmaceuticals", "FMCG / CPG", "Heavy Engineering",
    "Textiles", "Logistics", "Healthcare", "Retail"
]

TOOLS = [
    "Value Stream Mapping (VSM)", "5S & Workplace Org",
    "Hoshin Kanri", "Total Productive Maintenance (TPM)",
    "Six Sigma", "Lean", "Kanban"
]

REGIONS = [
    "India", "United States", "United Kingdom", "Germany",
    "France", "UAE", "Singapore", "Australia", "South Africa"
]

# --------------------------------------------------
# INTERNAL CASE STUDIES (STATIC)
# --------------------------------------------------
INTERNAL_CASES = [
    {
        "type": "case",
        "title": "Maruti Suzuki – Assembly Line VSM",
        "summary": "Value Stream Mapping reduced waste and rework across assembly lines.",
        "industry": "Automotive",
        "tool": "Value Stream Mapping (VSM)",
        "date": datetime(2024, 10, 15),
        "impact": "35% Cycle Time Reduction",
        "savings": "₹45 Cr"
    },
    {
        "type": "case",
        "title": "Apollo Hospitals – 5S Rollout",
        "summary": "5S deployment improved OT turnaround time and utilisation.",
        "industry": "Healthcare",
        "tool": "5S & Workplace Org",
        "date": datetime(2024, 9, 20),
        "impact": "27% OT Utilisation Increase",
        "savings": "₹12 Cr"
    }
]

# --------------------------------------------------
# FILE SCANNING
# --------------------------------------------------
def scan_folder(folder_path):
    files = []
    for root, _, filenames in os.walk(folder_path):
        for f in filenames:
            full_path = os.path.join(root, f)
            stat = os.stat(full_path)
            files.append({
                "title": f,
                "path": full_path,
                "ext": Path(f).suffix.lower(),
                "date": datetime.fromtimestamp(stat.st_mtime)
            })
    return files

# --------------------------------------------------
# TEXT EXTRACTION
# --------------------------------------------------
def extract_text(file):
    try:
        if file["ext"] == ".pdf":
            doc = fitz.open(file["path"])
            return " ".join(page.get_text() for page in doc)[:3000]

        if file["ext"] in [".ppt", ".pptx"]:
            prs = Presentation(file["path"])
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
            return text[:3000]
    except:
        return ""
    return ""

# --------------------------------------------------
# AI-STYLE SUMMARY + TAGS
# --------------------------------------------------
def ai_summary_and_tags(text):
    if not text:
        return "No preview available", ["Unclassified"]

    sentences = sent_tokenize(text)
    summary = " ".join(sentences[:2])

    tags = []
    t = text.lower()
    for k in ["vsm", "5s", "lean", "tpm", "six sigma", "kanban"]:
        if k in t:
            tags.append(k.upper())

    return summary, tags if tags else ["General"]

# --------------------------------------------------
# HEADER
# --------------------------------------------------
st.title("🚀 FABER NEXUS")
st.caption("Integrated Consulting Knowledge & Intelligence OS")
st.divider()

# --------------------------------------------------
# SIDEBAR
# --------------------------------------------------
with st.sidebar:
    industry = st.selectbox("Industry", INDUSTRIES)
    tool = st.selectbox("Framework", TOOLS)
    region = st.selectbox("Region", REGIONS)
    mode = st.radio("External Brain Mode", ["Curated (Guaranteed)", "Live (Best Effort)"])
    folder_path = st.text_input(
        "📁 Local Folder Path",
        value=r"D:\IIM Ranchi Sem 5\VLP Faber Infinite Consulting\Week 2 Updates"
    )

# --------------------------------------------------
# BUILD FILE-BASED CARDS (FIXED)
# --------------------------------------------------
FILE_CARDS = []

if not folder_path:
    st.warning("Please enter a folder path.")

elif not os.path.exists(folder_path):
    st.error("❌ Folder does not exist or is not accessible.")
    st.info("Run locally using: streamlit run app.py")

else:
    scanned_files = scan_folder(folder_path)
    st.success(f"📂 {len(scanned_files)} files detected")

    with st.expander("🔍 Debug: Files detected"):
        for f in scanned_files:
            st.write(f["title"])

    for f in scanned_files:
        text = extract_text(f)
        summary, tags = ai_summary_and_tags(text)

        FILE_CARDS.append({
            "type": "file",
            "title": f["title"],
            "summary": summary,
            "industry": "Internal",
            "tool": ", ".join(tags),
            "date": f["date"],
            "impact": "Document",
            "savings": "—",
            "path": f["path"]
        })

# --------------------------------------------------
# MERGE + SORT
# --------------------------------------------------
ALL_INTERNAL = sorted(
    INTERNAL_CASES + FILE_CARDS,
    key=lambda x: x["date"],
    reverse=True
)

# --------------------------------------------------
# TABS
# --------------------------------------------------
tab1, tab2, tab3 = st.tabs([
    "🧠 Internal Brain",
    "🌍 External Brain",
    "💰 ROI Simulator"
])

# --------------------------------------------------
# TAB 1 — INTERNAL BRAIN
# --------------------------------------------------
with tab1:
    st.subheader("📂 Internal Brain – Cases & Documents")

    f1, f2 = st.columns(2)
    ind_f = f1.selectbox("Filter by Industry", ["All"] + INDUSTRIES)
    tool_f = f2.selectbox("Filter by Tool", ["All"] + TOOLS)

    filtered = [
        c for c in ALL_INTERNAL
        if (ind_f == "All" or c["industry"] == ind_f or c["industry"] == "Internal")
        and (tool_f == "All" or tool_f in c["tool"])
    ]

    cols = st.columns(3)
    for i, c in enumerate(filtered):
        with cols[i % 3]:
            with st.container(border=True):
                st.markdown(f"### {c['title']}")
                st.caption(c["date"].strftime("%d %b %Y"))
                st.write(c["summary"])
                st.divider()
                st.markdown(f"**Impact:** {c['impact']}")
                st.markdown(f"**Savings:** {c['savings']}")
                st.caption(f"{c['industry']} • {c['tool']}")

                if c["type"] == "file":
                    if st.button("📂 Open File", key=c["path"]):
                        os.startfile(c["path"])

# --------------------------------------------------
# TAB 2 — EXTERNAL BRAIN (CURATED SAFE)
# --------------------------------------------------
with tab2:
    st.subheader("🌍 External Market Intelligence")

    benchmarks = [
        {
            "title": f"{industry} Operations Excellence Program",
            "summary": "Industry-wide programs typically deliver 20–30% cost reduction.",
            "savings": "₹25–50 Cr",
            "verified": False
        },
        {
            "title": f"{tool} Deployment – Global Case",
            "summary": "Framework-led transformations improve throughput by 25–40%.",
            "savings": "₹15–30 Cr",
            "verified": False
        }
    ]

    cols = st.columns(3)
    for i, r in enumerate(benchmarks):
        with cols[i % 3]:
            with st.container(border=True):
                badge = "indicative"
                st.markdown(
                    f"<span class='badge {badge}'>Indicative</span>",
                    unsafe_allow_html=True
                )
                st.markdown(f"### {r['title']}")
                st.write(r["summary"])
                st.divider()
                st.markdown(f"💰 {r['savings']}")

# --------------------------------------------------
# TAB 3 — ROI SIMULATOR
# --------------------------------------------------
with tab3:
    st.subheader("💰 ROI Simulator")

    revenue = st.number_input("Client Revenue (₹ Cr)", 100)
    ineff = st.slider("Inefficiency (%)", 5, 30, 15)
    fee = st.number_input("Consulting Fee (₹ Lakhs)", 25)

    savings = revenue * ineff / 100
    roi = (savings * 100 / fee) if fee > 0 else 0

    df = pd.DataFrame({
        "Category": ["Consulting Fee", "Projected Savings"],
        "₹ Crores": [fee / 100, savings]
    })

    fig = px.bar(df, x="Category", y="₹ Crores", text_auto=True)
    st.plotly_chart(fig, use_container_width=True)
    st.success(f"Projected ROI: {roi:.1f}x")

st.divider()
st.caption("Faber Infinite Consulting | Stable Local Build")
